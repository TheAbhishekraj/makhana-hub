import os
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas

# Directory configuration
EXPORT_DIR = r"C:\Users\HP\.gemini\antigravity\scratch\makhana-hub\exports"
os.makedirs(EXPORT_DIR, exist_ok=True)

# -------------------------------------------------------------
# 1. EXCEL EXPORTER
# -------------------------------------------------------------
def generate_excel():
    print("Generating Excel Financial Model...")
    wb = openpyxl.Workbook()
    
    # Styles
    font_family = "Segoe UI"
    font_title = Font(name=font_family, size=16, bold=True, color="FFFFFF")
    font_header = Font(name=font_family, size=11, bold=True, color="FFFFFF")
    font_bold = Font(name=font_family, size=10, bold=True)
    font_regular = Font(name=font_family, size=10)
    font_currency = Font(name=font_family, size=10)
    
    fill_forest = PatternFill(start_color="1B4332", end_color="1B4332", fill_type="solid")
    fill_mint = PatternFill(start_color="D8F3DC", end_color="D8F3DC", fill_type="solid")
    fill_gold = PatternFill(start_color="FEFAE0", end_color="FEFAE0", fill_type="solid")
    fill_zebra = PatternFill(start_color="F8F9FA", end_color="F8F9FA", fill_type="solid")
    
    thin_border_side = Side(border_style="thin", color="D3D3D3")
    thin_border = Border(left=thin_border_side, right=thin_border_side, top=thin_border_side, bottom=thin_border_side)
    double_bottom_border = Border(top=thin_border_side, bottom=Side(border_style="double", color="000000"))
    
    align_center = Alignment(horizontal="center", vertical="center")
    align_left = Alignment(horizontal="left", vertical="center")
    align_right = Alignment(horizontal="right", vertical="center")

    # SHEET 1: Cultivation Model
    ws1 = wb.active
    ws1.title = "Cultivation Model"
    ws1.views.sheetView[0].showGridLines = True
    
    # Title Block
    ws1.merge_cells("A1:D2")
    ws1["A1"] = "1-Hectare Field-Based Cultivation Model (Swarna Vaidehi)"
    ws1["A1"].font = font_title
    ws1["A1"].fill = fill_forest
    ws1["A1"].alignment = align_center
    ws1.row_dimensions[1].height = 20
    ws1.row_dimensions[2].height = 20
    
    # Headers
    headers1 = ["Operational / Input Item", "Quantity / Specifications", "Unit Cost (INR)", "Total Cost (INR)"]
    ws1.append([]) # row 3 empty
    ws1.row_dimensions[4].height = 24
    
    for col_idx, header in enumerate(headers1, start=1):
        cell = ws1.cell(row=4, column=col_idx, value=header)
        cell.font = font_header
        cell.fill = fill_forest
        cell.alignment = align_center
        cell.border = thin_border
        
    opex_items = [
        ("Land Lease (1 Year)", "Clayey soil bunded field", 25000),
        ("Field Preparation & Perimeter Bunding", "Mechanical earth-moving + leveling (30-45 cm)", 15000),
        ("Nursery Setup & Certified Seeds", "20 kg Euryale certified seed (Ponds require 90 kg)", 8000),
        ("Nutrition Regimen (NPK + Micro-nutrients)", "NPK @ 100:60:40 + S(20kg), Zn(5kg), Mg(5kg), B(5kg)", 12000),
        ("IPM Protection & Bio-agents Sourcing", "Metarhizium anisopliae + Chlorantraniliprole/Neem", 5000),
        ("Field Labor mandays", "75 Mandays @ Rs. 600 (Transplanting/Harvesting)", 45000),
        ("Irrigation Water & Pump Fuel", "Diesel replenishment during summer dry spells", 10000)
    ]
    
    for idx, (item, spec, cost) in enumerate(opex_items, start=5):
        ws1.append([item, spec, cost, f"=C{idx}"])
        row = ws1.row_dimensions[idx]
        row.height = 20
        ws1[f"A{idx}"].font = font_regular
        ws1[f"B{idx}"].font = font_regular
        ws1[f"C{idx}"].font = font_currency
        ws1[f"C{idx}"].number_format = '₹#,##0'
        ws1[f"D{idx}"].font = font_bold
        ws1[f"D{idx}"].number_format = '₹#,##0'
        
        for c in range(1, 5):
            ws1.cell(row=idx, column=c).border = thin_border
            if idx % 2 == 0:
                ws1.cell(row=idx, column=c).fill = fill_zebra
                
    # Total Opex Row
    tot_row = len(opex_items) + 5
    ws1.cell(row=tot_row, column=1, value="Total Cultivation OPEX").font = font_bold
    ws1.cell(row=tot_row, column=4, value=f"=SUM(D5:D{tot_row-1})").font = font_bold
    ws1.cell(row=tot_row, column=4).number_format = '₹#,##0'
    ws1.cell(row=tot_row, column=1).alignment = align_left
    ws1.cell(row=tot_row, column=4).alignment = align_right
    for c in range(1, 5):
        ws1.cell(row=tot_row, column=c).border = double_bottom_border
        ws1.cell(row=tot_row, column=c).fill = fill_mint

    # Revenue Section
    rev_start = tot_row + 3
    ws1.merge_cells(start_row=rev_start, start_column=1, end_row=rev_start+1, end_column=4)
    ws1.cell(row=rev_start, column=1, value="Operational Yield & Revenue Generation").font = font_title
    ws1.cell(row=rev_start, column=1).fill = fill_forest
    ws1.cell(row=rev_start, column=1).alignment = align_center
    
    ws1.cell(row=rev_start+3, column=1, value="Output Product").font = font_header
    ws1.cell(row=rev_start+3, column=1).fill = fill_forest
    ws1.cell(row=rev_start+3, column=2, value="Yield Volume (Kg)").font = font_header
    ws1.cell(row=rev_start+3, column=2).fill = fill_forest
    ws1.cell(row=rev_start+3, column=3, value="Rate per Kg (INR)").font = font_header
    ws1.cell(row=rev_start+3, column=3).fill = fill_forest
    ws1.cell(row=rev_start+3, column=4, value="Total Revenue (INR)").font = font_header
    ws1.cell(row=rev_start+3, column=4).fill = fill_forest
    
    ws1.cell(row=rev_start+4, column=1, value="Raw Makhana Seeds (Swarna Vaidehi)").font = font_regular
    ws1.cell(row=rev_start+4, column=2, value=2800).font = font_regular
    ws1.cell(row=rev_start+4, column=3, value=150).font = font_currency
    ws1.cell(row=rev_start+4, column=4, value=f"=B{rev_start+4}*C{rev_start+4}").font = font_bold
    ws1.cell(row=rev_start+4, column=4).number_format = '₹#,##0'
    
    for c in range(1, 5):
        ws1.cell(row=rev_start+3, column=c).border = thin_border
        ws1.cell(row=rev_start+3, column=c).alignment = align_center
        ws1.cell(row=rev_start+4, column=c).border = thin_border
        
    ws1.cell(row=rev_start+5, column=1, value="Net Cultivator Profit per Hectare").font = font_bold
    ws1.cell(row=rev_start+5, column=4, value=f"=D{rev_start+4}-D{tot_row}").font = font_bold
    ws1.cell(row=rev_start+5, column=4).number_format = '₹#,##0'
    ws1.cell(row=rev_start+5, column=1).alignment = align_left
    ws1.cell(row=rev_start+5, column=4).alignment = align_right
    for c in range(1, 5):
        ws1.cell(row=rev_start+5, column=c).border = double_bottom_border
        ws1.cell(row=rev_start+5, column=c).fill = fill_mint

    # SHEET 2: Processing Ledger
    ws2 = wb.create_sheet(title="Processing Ledger Y1")
    ws2.views.sheetView[0].showGridLines = True
    
    ws2.merge_cells("A1:E2")
    ws2["A1"] = "Year-1 Bootstrapped Processing Ledger (30 Tons Raw Seeds)"
    ws2["A1"].font = font_title
    ws2["A1"].fill = fill_forest
    ws2["A1"].alignment = align_center
    
    headers2 = ["Operational Component", "Description / Details", "Volume / Metric", "Rate / Price", "Amount (INR)"]
    ws2.append([])
    ws2.row_dimensions[4].height = 24
    for col_idx, header in enumerate(headers2, start=1):
        cell = ws2.cell(row=4, column=col_idx, value=header)
        cell.font = font_header
        cell.fill = fill_forest
        cell.alignment = align_center
        cell.border = thin_border
        
    proc_rows = [
        ("Raw Seed Procurement", "30 Tons direct buyback at farm gate", 30000, 130, "=C5*D5"),
        ("Primary Drying & Grading", "Sorting into size classes at centers", 150, 800, "=C6*D6"),
        ("Logistics & Transport", "Hired small transport to clusters", 1, 150000, "=C7*D7"),
        ("Manual Popping Job-Work", "39% weight recovery popping fee", 11700, 35, "=C8*D8"),
        ("Packaging & Branding", "Nitrogen flushed custom pouches", 46800, 2.56, "=C9*D9"),
        ("Hub Overhead / Rent", "Warehouse rent & general overheads", 12, 12500, "=C10*D10")
    ]
    
    for idx, (comp, desc, vol, rate, amt) in enumerate(proc_rows, start=5):
        ws2.append([comp, desc, vol, rate, amt])
        ws2.row_dimensions[idx].height = 20
        ws2[f"A{idx}"].font = font_regular
        ws2[f"B{idx}"].font = font_regular
        ws2[f"C{idx}"].font = font_regular
        ws2[f"D{idx}"].font = font_regular
        ws2[f"E{idx}"].font = font_bold
        ws2[f"E{idx}"].number_format = '₹#,##0'
        for c in range(1, 6):
            ws2.cell(row=idx, column=c).border = thin_border
            if idx % 2 == 0:
                ws2.cell(row=idx, column=c).fill = fill_zebra
                
    tot_proc_row = len(proc_rows) + 5
    ws2.cell(row=tot_proc_row, column=1, value="Total Processing & Sourcing OPEX").font = font_bold
    ws2.cell(row=tot_proc_row, column=5, value=f"=SUM(E5:E{tot_proc_row-1})").font = font_bold
    ws2.cell(row=tot_proc_row, column=5).number_format = '₹#,##0'
    for c in range(1, 6):
        ws2.cell(row=tot_proc_row, column=c).border = double_bottom_border
        ws2.cell(row=tot_proc_row, column=c).fill = fill_mint

    # Processing Revenue
    rev_start_p = tot_proc_row + 3
    ws2.merge_cells(start_row=rev_start_p, start_column=1, end_row=rev_start_p+1, end_column=5)
    ws2.cell(row=rev_start_p, column=1, value="Hub Popping Revenue & Yield Outputs").font = font_title
    ws2.cell(row=rev_start_p, column=1).fill = fill_forest
    ws2.cell(row=rev_start_p, column=1).alignment = align_center
    
    ws2.cell(row=rev_start_p+3, column=1, value="Popped Product Grade").font = font_header
    ws2.cell(row=rev_start_p+3, column=1).fill = fill_forest
    ws2.cell(row=rev_start_p+3, column=2, value="Output specifications").font = font_header
    ws2.cell(row=rev_start_p+3, column=2).fill = fill_forest
    ws2.cell(row=rev_start_p+3, column=3, value="Quantity Sold (Kg)").font = font_header
    ws2.cell(row=rev_start_p+3, column=3).fill = fill_forest
    ws2.cell(row=rev_start_p+3, column=4, value="Wholesale Price / Kg").font = font_header
    ws2.cell(row=rev_start_p+3, column=4).fill = fill_forest
    ws2.cell(row=rev_start_p+3, column=5, value="Total Revenue").font = font_header
    ws2.cell(row=rev_start_p+3, column=5).fill = fill_forest
    
    rev_rows = [
        ("Premium Popped (Grade-A)", "Fully popped premium kernels (> 1.2 cm)", 8000, 550, f"=C{rev_start_p+4}*D{rev_start_p+4}"),
        ("Standard Popped (Grade-B)", "Perfect popping standard sizing (0.9-1.2 cm)", 3000, 400, f"=C{rev_start_p+5}*D{rev_start_p+5}"),
        ("Broken Makhana / Phulla", "Secondary sorting sold for weaning flour", 700, 200, f"=C{rev_start_p+6}*D{rev_start_p+6}"),
        ("Makhana Shell Waste", "Mulching & industrial biomass fuel waste", 18000, 3, f"=C{rev_start_p+7}*D{rev_start_p+7}")
    ]
    
    for idx, (gr, spec, q, r, formula) in enumerate(rev_rows, start=rev_start_p+4):
        ws2.append([gr, spec, q, r, formula])
        ws2.row_dimensions[idx].height = 20
        ws2[f"A{idx}"].font = font_regular
        ws2[f"B{idx}"].font = font_regular
        ws2[f"C{idx}"].font = font_regular
        ws2[f"D{idx}"].font = font_regular
        ws2[f"E{idx}"].font = font_bold
        ws2[f"E{idx}"].number_format = '₹#,##0'
        for c in range(1, 6):
            ws2.cell(row=idx, column=c).border = thin_border
            if idx % 2 == 0:
                ws2.cell(row=idx, column=c).fill = fill_zebra
                
    tot_rev_p = rev_start_p + 8
    ws2.cell(row=tot_rev_p, column=1, value="Total Hub Revenue").font = font_bold
    ws2.cell(row=tot_rev_p, column=5, value=f"=SUM(E{rev_start_p+4}:E{tot_rev_p-1})").font = font_bold
    ws2.cell(row=tot_rev_p, column=5).number_format = '₹#,##0'
    for c in range(1, 6):
        ws2.cell(row=tot_rev_p, column=c).border = double_bottom_border
        ws2.cell(row=tot_rev_p, column=c).fill = fill_mint
        
    net_p_row = tot_rev_p + 1
    ws2.cell(row=net_p_row, column=1, value="Net Hub Processing Profit (Year 1)").font = font_bold
    ws2.cell(row=net_p_row, column=5, value=f"=E{tot_rev_p}-E{tot_proc_row}").font = font_bold
    ws2.cell(row=net_p_row, column=5).number_format = '₹#,##0'
    for c in range(1, 6):
        ws2.cell(row=net_p_row, column=c).border = double_bottom_border
        ws2.cell(row=net_p_row, column=c).fill = fill_gold

    # SHEET 3: 5-Ton Mechanized Scale
    ws3 = wb.create_sheet(title="5-Ton Scale CAPEX OPEX")
    ws3.views.sheetView[0].showGridLines = True
    
    ws3.merge_cells("A1:C2")
    ws3["A1"] = "5-Ton Fully Mechanized Scaling Phase (Year 3 Onwards)"
    ws3["A1"].font = font_title
    ws3["A1"].fill = fill_forest
    ws3["A1"].alignment = align_center
    
    # CAPEX Columns
    ws3.cell(row=4, column=1, value="Initial CAPEX Machinery & Infrastructure").font = font_header
    ws3.cell(row=4, column=1).fill = fill_forest
    ws3.cell(row=4, column=2, value="Technical Specifications").font = font_header
    ws3.cell(row=4, column=2).fill = fill_forest
    ws3.cell(row=4, column=3, value="Cost (INR)").font = font_header
    ws3.cell(row=4, column=3).fill = fill_forest
    
    # Exact values from Google Drive File 1
    capex_items = [
        ("Secondary Roasting & Popping Machine", "Impact popping at 270 °C (25-30 kg/h)", 800000),
        ("Primary Roasting Machine", "Roasts seeds at 270 °C for 3-5 mins (100 kg/h)", 140000),
        ("Rotary Nested Grader", "Separates raw seeds into 7 size classes (180 kg/h)", 250000),
        ("Makhana Seed Washer", "Cleans raw seeds by friction (120-180 kg/h)", 160000),
        ("Makhana Seed Grader", "Sand suction slurry sorting unit (0.5 t/h)", 130000),
        ("Cabinet Seed Dryer", "Temp range 30-220 °C (20-30 kg/h)", 200000),
        ("Popped Makhana Grader", "Sorts finished puffs into 7 size classes (100 kg/h)", 200000),
        ("Form, Fill & Seal Packing Machine", "Multi-head weigher automatic packaging (6-8 t/d)", 800000),
        ("Makhana Polishing Unit", "Polishes pops for smooth texture", 100000),
        ("Warehouse Civil Works & Drying Slab", "Concrete slab and dust-free sorting floor", 1500000),
        ("Solar Rooftop Micro-Grid", "15kW solar array with 10kW battery backup grid", 500000)
    ]
    
    for idx, (item, spec, val) in enumerate(capex_items, start=5):
        ws3.append([item, spec, val])
        ws3.row_dimensions[idx].height = 20
        ws3[f"A{idx}"].font = font_regular
        ws3[f"B{idx}"].font = font_regular
        ws3[f"C{idx}"].font = font_bold
        ws3[f"C{idx}"].number_format = '₹#,##0'
        for c in range(1, 4):
            ws3.cell(row=idx, column=c).border = thin_border
            if idx % 2 == 0:
                ws3.cell(row=idx, column=c).fill = fill_zebra
                
    tot_capex_row = len(capex_items) + 5
    ws3.cell(row=tot_capex_row, column=1, value="Total Capital Required (CAPEX)").font = font_bold
    ws3.cell(row=tot_capex_row, column=3, value=f"=SUM(C5:C{tot_capex_row-1})").font = font_bold
    ws3.cell(row=tot_capex_row, column=3).number_format = '₹#,##0'
    for c in range(1, 4):
        ws3.cell(row=tot_capex_row, column=c).border = double_bottom_border
        ws3.cell(row=tot_capex_row, column=c).fill = fill_mint

    # Monthly OPEX Columns (Side-by-side or below)
    op_start = tot_capex_row + 3
    ws3.merge_cells(start_row=op_start, start_column=1, end_row=op_start+1, end_column=3)
    ws3.cell(row=op_start, column=1, value="Monthly Operational Expenditure (OPEX @ 25 Tons)").font = font_title
    ws3.cell(row=op_start, column=1).fill = fill_forest
    ws3.cell(row=op_start, column=1).alignment = align_center
    
    ws3.cell(row=op_start+3, column=1, value="Monthly OPEX Input Category").font = font_header
    ws3.cell(row=op_start+3, column=1).fill = fill_forest
    ws3.cell(row=op_start+3, column=2, value="Activity Detail / Specification").font = font_header
    ws3.cell(row=op_start+3, column=2).fill = fill_forest
    ws3.cell(row=op_start+3, column=3, value="Cost per Month (INR)").font = font_header
    ws3.cell(row=op_start+3, column=3).fill = fill_forest
    
    opex_m_items = [
        ("Raw Material Sourcing", "25,000 kg seeds @ average Rs. 130/kg", 3250000),
        ("Plant Electricity & Roaster Fuel (LDO)", "Combustion burners and automated air systems", 45000),
        ("Plant Factory Labor", "1 Supervisor, 2 operators, 6 sorting hands", 120000),
        ("Logistics & National Freight", "Seed haulage + outward distribution logistics", 100000),
        ("Factory Maintenance & Depreciation", "Equipment maintenance and backup fund reserves", 35000),
        ("Compliance, Marketing & FSSAI", "Brand packaging + export safety compliance (FSSAI/EU)", 50000)
    ]
    
    for idx, (item, spec, val) in enumerate(opex_m_items, start=op_start+4):
        ws3.append([item, spec, val])
        ws3.row_dimensions[idx].height = 20
        ws3[f"A{idx}"].font = font_regular
        ws3[f"B{idx}"].font = font_regular
        ws3[f"C{idx}"].font = font_bold
        ws3[f"C{idx}"].number_format = '₹#,##0'
        for c in range(1, 4):
            ws3.cell(row=idx, column=c).border = thin_border
            if idx % 2 == 0:
                ws3.cell(row=idx, column=c).fill = fill_zebra
                
    tot_opexm_row = op_start + len(opex_m_items) + 4
    ws3.cell(row=tot_opexm_row, column=1, value="Total Monthly Scaling OPEX").font = font_bold
    ws3.cell(row=tot_opexm_row, column=3, value=f"=SUM(C{op_start+4}:C{tot_opexm_row-1})").font = font_bold
    ws3.cell(row=tot_opexm_row, column=3).number_format = '₹#,##0'
    for c in range(1, 4):
        ws3.cell(row=tot_opexm_row, column=c).border = double_bottom_border
        ws3.cell(row=tot_opexm_row, column=c).fill = fill_mint

    # SHEET 4: IFS vs Monoculture
    ws4 = wb.create_sheet(title="IFS Comparison")
    ws4.views.sheetView[0].showGridLines = True
    
    ws4.merge_cells("A1:C2")
    ws4["A1"] = "Integrated Farming System (IFS) vs Monoculture Comparison"
    ws4["A1"].font = font_title
    ws4["A1"].fill = fill_forest
    ws4["A1"].alignment = align_center
    
    headers4 = ["Comparison Metric / Yield Component", "Traditional Field Monoculture", "Makhana + Carps + Chestnut IFS"]
    ws4.append([])
    ws4.row_dimensions[4].height = 24
    for col_idx, header in enumerate(headers4, start=1):
        cell = ws4.cell(row=4, column=col_idx, value=header)
        cell.font = font_header
        cell.fill = fill_forest
        cell.alignment = align_center
        cell.border = thin_border
        
    ifs_rows = [
        ("Water Depth Maintenance", "30 - 45 cm (Shallow)", "50 - 60 cm (For fish movement)"),
        ("Primary Makhana Yield (Kg/Ha)", 2800, 2600),
        ("Carps Fish Yield (Kg/Ha)", 0, 1200),
        ("Water Chestnut Yield (Kg/Ha)", 0, 3000),
        ("Total Operational Cost (INR)", 120000, 180000),
        ("Total Gross Revenue (INR)", 420000, 730000),
        ("Net Farm Profit per Hectare (INR)", "=B10-B9", "=C10-C9")
    ]
    
    for idx, (met, mono, ifs) in enumerate(ifs_rows, start=5):
        ws4.append([met, mono, ifs])
        ws4.row_dimensions[idx].height = 20
        ws4[f"A{idx}"].font = font_regular
        ws4[f"B{idx}"].font = font_bold if str(mono).startswith("=") else font_regular
        ws4[f"C{idx}"].font = font_bold
        
        if str(mono).isdigit() or str(mono).startswith("="):
            ws4[f"B{idx}"].number_format = '₹#,##0' if "Cost" in met or "Revenue" in met or "Profit" in met else '#,##0'
        if str(ifs).isdigit() or str(ifs).startswith("="):
            ws4[f"C{idx}"].number_format = '₹#,##0' if "Cost" in met or "Revenue" in met or "Profit" in met else '#,##0'
            
        for c in range(1, 4):
            ws4.cell(row=idx, column=c).border = thin_border
            if idx == 11:
                ws4.cell(row=idx, column=c).fill = fill_mint
                ws4.cell(row=idx, column=c).border = double_bottom_border
                
    # Auto-adjust column widths across all sheets
    for sheet in wb.worksheets:
        for col in sheet.columns:
            max_len = 0
            col_letter = get_column_letter(col[0].column)
            for cell in col:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            sheet.column_dimensions[col_letter].width = max(max_len + 3, 12)
            
    # Specific tuning for large columns
    ws1.column_dimensions["A"].width = 38
    ws1.column_dimensions["B"].width = 38
    ws2.column_dimensions["A"].width = 28
    ws2.column_dimensions["B"].width = 38
    ws3.column_dimensions["A"].width = 35
    ws3.column_dimensions["B"].width = 45
    ws4.column_dimensions["A"].width = 35
    ws4.column_dimensions["B"].width = 28
    ws4.column_dimensions["C"].width = 28

    wb.save(os.path.join(EXPORT_DIR, "makhana_hub_financial_model.xlsx"))
    print("Excel Financial Model saved successfully!")

# -------------------------------------------------------------
# 2. POWERPOINT EXPORTER
# -------------------------------------------------------------
def generate_ppt():
    print("Generating Slide Deck...")
    prs = Presentation()
    
    # Standard 16:9 layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 12 slides data with exact facts
    slides_data = [
        {
            "slide": 1,
            "label": "Introduction",
            "title": "The Swarna Vaidehi Revolution",
            "subtitle": "Building a World-Class Makhana Processing Hub in North Bihar",
            "bullets": [
                "Strategic Vision: Scaling Euryale ferox (Makhana) sector of Bihar.",
                "Agronomy Foundation: Shallow-field elite variety developed by ICAR-NRC Makhana.",
                "CVRC Notified Variety: Released on 15.11.2013, notified on 30.07.2022.",
                "Digital Link: Fully integrated with the Bihar Bazar platform (bihar-bazaar.vercel.app)."
            ],
            "notes": "Good morning, partners. Today we present an opportunity to organize the GI-tagged makhana sector of Bihar. By combining the field-cultivation technology of ICAR-NRC Darbhanga with a modern, high-value value-add supply chain, we are unlocking massive margins in a premium global superfood segment."
        },
        {
            "slide": 2,
            "label": "The Core Problem",
            "title": "Primitive Harvesting & Inefficient Supply Chains",
            "subtitle": "Severe physical labor bottlenecks and structural value loss",
            "bullets": [
                "Extreme Physical Danger: Traditional harvesting requires deep pond diving (1.5-2.0m depth).",
                "High Waste Sourcing: Traditional seed rate is a high 90 kg/ha with massive post-harvest mold.",
                "Middleman Exploitation: Traditional 'Arhatiyas' capture up to 60% of farm-gate value.",
                "Manual Popping Bottlenecks: Primitive manual popping is slow and restricts volume, limiting value-add scaling."
            ],
            "notes": "Traditional pond harvesting is physically dangerous, causing severe respiratory and skin infections for Mallah divers. Unorganized trade strips local farmers of their fair share, while poor post-harvest handling leads to degraded seed quality and poor popping expansion."
        },
        {
            "slide": 3,
            "label": "The Scalable Solution",
            "title": "The Integrated Makhana Hub Ecosystem",
            "subtitle": "Re-engineering harvesting safety and product logistics",
            "bullets": [
                "Field-Based Transition: Safe, shallow field bunding (30-45 cm depth) replaces deep ponds.",
                "Slashed Sowing Seed Rate: Seed rate cut from 90 kg/ha down to 20 kg/ha via nursery setups.",
                "Decentralized Popping Network: Zero machinery CAPEX in Year 1 using traditional popping clusters.",
                "Value-Added Nutrition: Grinding phulla into weaning baby food and making premium gluten-free cookies."
            ],
            "notes": "We are introducing a dual solution: we move cultivation into shallow agricultural fields, saving labor and water, and we set up decentralized primary processing centers. In Year 1, we bootstrap operations by outsourcing popping to traditional manual clusters before moving to fully automated lines."
        },
        {
            "slide": 4,
            "label": "Scientific Superiority",
            "title": "Agronomy of 'Swarna Vaidehi' Variety",
            "subtitle": "Breakthrough field performance and water conservation",
            "bullets": [
                "Unprecedented Yield: 3.0 - 3.5 Tons / Ha (Average 25 q/ha) compared to pond average of 1.4 Tons.",
                "Nursery Bed Sowing: Nov-March with nursery area of 500 m²/ha. Transplanting in Feb at 1.20m spacing.",
                "Nutrient Regimen: 100:60:40 NPK kg/ha + 10 t/ha FYM + 5 kg/ha Zn, Mg, B micronutrient setup.",
                "Foliar Micronutrients: Sprays of 0.1% CuSO4 solution significantly boost shell thickness."
            ],
            "notes": "Developed by ICAR-NRC Makhana, Swarna Vaidehi doubles traditional yield. The optimized shallow field-based environment requires 80% less water. Precision secondary nutrients (S, Zn, Mg, B) ensure thick, durable seed shells, yielding high popping volumes."
        },
        {
            "slide": 5,
            "label": "Procurement & Quality Control",
            "title": "Distributed Sourcing & Grading Logistics",
            "subtitle": "Eliminating middleman leakages through digital collection points",
            "bullets": [
                "Strategic Sourcing Nodes: Darbhanga (Center), Madhubani (North), and Purnia (East).",
                "HDPE Pre-Drying: Washing and sun-drying on 250 GSM sheets to reduce moisture from 28% to 12-14%.",
                "Rotary Nested Grader: Size sorting into 5 grades at 180 kg/h to ensure roasting at 270 °C without burning.",
                "Manual Popping Cluster: Skilled traditional popping families compensated at ₹35/kg popped makhana."
            ],
            "notes": "By setting up direct procurement centers in Darbhanga, Madhubani, and Purnia, we control raw seed quality at the farm-gate. Uniform mechanical sizing prevents roasted seeds from burning, ensuring a premium puffed product."
        },
        {
            "slide": 6,
            "label": "Product Innovation & Chemistry",
            "title": "Premium Functional Superfoods",
            "subtitle": "Re-engineering Euryale ferox digestible starchy perisperm core",
            "bullets": [
                "Makhana Ragi Cookies: Rich in calcium (160mg/100g) and amino acids (Glutamine, Arginine) - 58% margin.",
                "Weaning Baby Food (Shishu Aahar): Hypoallergenic formula with optimized 1.2:1 Ca:P ratio (62% margin).",
                "Low-GI Starch Pasta: Extruded cold-gelatinized high-amylose starch pasta (52% margin).",
                "Perisperm Value: Capitalizing on Euryale ferox starchy core's low fat (0.1%) and high protein (10-12%)."
            ],
            "notes": "Instead of selling raw commodities, we process broken makhana into specialized high-margin products. Micronized flour cookies and low-GI pasta tap into the massive global wellness boom, delivering net margins exceeding 50%."
        },
        {
            "slide": 7,
            "label": "Financial DPR Highlights",
            "title": "Year-1 Zero Machinery Profitability",
            "subtitle": "A lean, bootstrapped operational ledger delivering immediate returns",
            "bullets": [
                "1-Hectare Cultivation: OPEX ₹1.2 Lakhs vs Revenue ₹4.2 Lakhs (₹3,00,000 Net Profit).",
                "30-Ton Processing Hub: OPEX ₹48.5 Lakhs vs Revenue ₹57.9 Lakhs (₹9,44,500 Net Profit).",
                "Consolidated Year-1 Profit: ₹12,44,500 achieved securely.",
                "Zero Initial CAPEX: Outsourcing popping to skilled manual clusters keeps startup costs low."
            ],
            "notes": "Our financial model is highly resilient. In Year 1, we avoid high equipment costs by leveraging traditional manual popping clusters. A 30-ton processing volume alone generates over ₹9.4 Lakhs in net profit, which when combined with our 1-Ha self-cultivation delivers ₹12.4 Lakhs in aggregate gains."
        },
        {
            "slide": 8,
            "label": "The 5-Year Scaling Roadmap",
            "title": "Phased Capital Expansion Blueprint",
            "subtitle": "Leaping from manual artisan networks to high-throughput automation",
            "bullets": [
                "Year 1-2: Establish manual processing hubs, secure organic certificates, and expand sourcing.",
                "Year 3: Raise institutional funding to build a 5-ton/day fully mechanized popping plant.",
                "Year 4: Launch direct-to-consumer brand lines and expand product lines.",
                "Year 5: Secure international export markets (US, EU, GCC) and achieve circular zero-waste."
            ],
            "notes": "We start lean. Once our direct farm sourcing and distribution lines are fully validated in Years 1 and 2, we will leverage bank debt and subsidies to install an automated 5-ton/day plant. This will cut processing costs by 45% and scale output dramatically."
        },
        {
            "slide": 9,
            "label": "Risk & Climate Resilience",
            "title": "Climate-Adaptive Agricultural Infrastructure",
            "subtitle": "Engineering field stability against monsoons and dry seasons",
            "bullets": [
                "Clay-Lined Perimeter Bunding: Prevents lateral water seepage, maintaining soil moisture.",
                "Solar Drip & Fill Pumps: Automated irrigation control during summer droughts.",
                "Botanical Protection: Biocontrol agents and garlic-chilli sprays replace synthetic chemistry.",
                "Financial Protection: Pre-season price agreements and multi-peril crop insurance."
            ],
            "notes": "Climate risk is mitigated through off-grid solar water pumps and high perimeter bunds lined with heavy clay. Integrated biological pest control (Trichoderma and botanical sprays) protects the crop while keeping it chemical-residue free."
        },
        {
            "slide": 10,
            "label": "Digital Integration",
            "title": "Eliminating Middlemen via Bihar Bazar",
            "subtitle": "Closing localized farmer literacy gaps and structural supply-chain asymmetry",
            "bullets": [
                "Direct Web Integration: Listing on the Bihar Bazar platform (bihar-bazaar.vercel.app).",
                "Voice & SMS Portals: Direct mobile messaging bypasses local digital and literacy gaps.",
                "Transparent Price Engine: Real-time grade-wise pricing determined by automated sieve sorting.",
                "Traceability: QR-code tracking from Bihar fields directly to premium retail shelves."
            ],
            "notes": "By integrating with the Bihar Bazar digital platform, we bring transparency to farm-gate transactions. Farmers use voice and SMS to lock in purchase prices, bypassing traditional middleman networks entirely."
        },
        {
            "slide": 11,
            "label": "Government Support & Schemes",
            "title": "PMFME & Mudra Capital Alignment",
            "subtitle": "Leveraging central subsidies and interest subventions",
            "bullets": [
                "PMFME Capital Subsidy: Claiming 35% capital grant (up to ₹10 Lakhs limit).",
                "Mudra Bank Funding: Collateral-free low-interest working capital credit lines.",
                "Agriculture Infrastructure Fund: 3% interest subvention on long-term infrastructure loans.",
                "State Policy Rebates: State industrial incentives including power exemptions and tax benefits."
            ],
            "notes": "We are fully aligned with the central government's food processing push. We leverage PMFME capital grants and central interest subventions to lower our borrowing costs, accelerating shareholder equity returns."
        },
        {
            "slide": 12,
            "label": "The Investment Ask & Returns",
            "title": "Capital Proposition & Socio-Economic Impact",
            "subtitle": "High financial yields combined with sustainable rural development",
            "bullets": [
                "The Ask: ₹50 Lakhs institutional funding (mix of 70% bank debt and 30% founder equity).",
                "Key Returns: Projected 5-Year Financial IRR of 38.6% with a 2.4-year Payback Period.",
                "Social Impact: Steady wages for 120+ rural women popping artisans and 250+ farmers.",
                "Exit Options: Strategic acquisition by fast-growing national agritech brands."
            ],
            "notes": "We are seeking ₹50 Lakhs to transition into our mechanized scaling phase. We offer a high project IRR of 38.6% and a quick 2.4-year payback, coupled with a deep social impact story. Join us in scaling Bihar's white gold."
        }
    ]
    
    # Forest Green theme colors
    color_bg = RGBColor(11, 37, 26)       # #0B251A
    color_title = RGBColor(255, 255, 255) # #FFFFFF
    color_gold = RGBColor(245, 203, 92)   # #F5CB5C
    color_bullet = RGBColor(220, 225, 222) # #DCE1DE
    
    for s in slides_data:
        # Create blank slide
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set solid dark forest green background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_bg
        
        # Add Slide Label top left
        txBox_lbl = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.5))
        tf_lbl = txBox_lbl.text_frame
        p_lbl = tf_lbl.paragraphs[0]
        p_lbl.text = f"SLIDE {s['slide']} : {s['label'].upper()}"
        p_lbl.font.size = Pt(11)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = color_gold
        p_lbl.font.name = "Segoe UI"
        
        # Add Slide Title
        txBox_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(0.8))
        tf_title = txBox_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = s["title"]
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = color_title
        p_title.font.name = "Segoe UI"
        
        # Add Subtitle
        txBox_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(0.5))
        tf_sub = txBox_sub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = s["subtitle"]
        p_sub.font.size = Pt(16)
        p_sub.font.italic = True
        p_sub.font.color.rgb = color_gold
        p_sub.font.name = "Segoe UI"
        
        # Add Bullet Points
        txBox_bullets = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(11.5), Inches(4.5))
        tf_bullets = txBox_bullets.text_frame
        tf_bullets.word_wrap = True
        
        for idx, bullet in enumerate(s["bullets"]):
            if idx == 0:
                p_b = tf_bullets.paragraphs[0]
            else:
                p_b = tf_bullets.add_paragraph()
            p_b.text = f"•  {bullet}"
            p_b.space_after = Pt(16)
            p_b.font.size = Pt(18)
            p_b.font.color.rgb = color_bullet
            p_b.font.name = "Segoe UI"
            
        # Set Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = s["notes"]
        
    prs.save(os.path.join(EXPORT_DIR, "makhana_hub_pitch_deck.pptx"))
    print("PowerPoint Pitch Deck saved successfully!")

# -------------------------------------------------------------
# 3. PDF REPORT EXPORTER (REPORTLAB)
# -------------------------------------------------------------
class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        canvas.Canvas.__init__(self, *args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_number(num_pages)
            canvas.Canvas.showPage(self)
        canvas.Canvas.save(self)

    def draw_page_number(self, page_count):
        if self._pageNumber == 1:
            return  # Skip page numbering on cover page
        self.saveState()
        self.setFont("Helvetica", 9)
        self.setFillColor(colors.HexColor("#475569"))
        
        # Header
        self.drawString(54, 750, "Swarna Vaidehi Makhana Hub — 5-Year Strategic Roadmap & DPR")
        self.setStrokeColor(colors.HexColor("#E2E8F0"))
        self.setLineWidth(0.5)
        self.line(54, 742, 558, 742)
        
        # Footer
        page_text = f"Page {self._pageNumber} of {page_count}"
        self.drawRightString(558, 40, page_text)
        self.drawString(54, 40, "CONFIDENTIAL — FOR BANKING AND INSTITUTIONAL EVALUATION ONLY")
        self.line(54, 52, 558, 52)
        self.restoreState()

def generate_pdf():
    print("Generating PDF Report...")
    pdf_path = os.path.join(EXPORT_DIR, "makhana_hub_strategic_report.pdf")
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=letter,
        rightMargin=54,
        leftMargin=54,
        topMargin=72,
        bottomMargin=72
    )
    
    styles = getSampleStyleSheet()
    
    # Custom styled paragraphs
    title_style = ParagraphStyle(
        'CoverTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=32,
        leading=38,
        textColor=colors.HexColor("#1B4332"),
        spaceAfter=15,
        alignment=1
    )
    
    subtitle_style = ParagraphStyle(
        'CoverSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-Oblique',
        fontSize=14,
        leading=18,
        textColor=colors.HexColor("#606C38"),
        spaceAfter=30,
        alignment=1
    )
    
    meta_style = ParagraphStyle(
        'CoverMeta',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#475569"),
        alignment=1
    )
    
    h1_style = ParagraphStyle(
        'Header1',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=20,
        leading=24,
        textColor=colors.HexColor("#1B4332"),
        spaceBefore=18,
        spaceAfter=10,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'Header2',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=colors.HexColor("#2D6A4F"),
        spaceBefore=12,
        spaceAfter=6,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'BodyTextCustom',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#334155"),
        spaceAfter=8
    )
    
    bullet_style = ParagraphStyle(
        'BulletCustom',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#334155"),
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=4
    )
    
    table_header_style = ParagraphStyle(
        'TableHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9,
        leading=11,
        textColor=colors.white
    )
    
    table_body_style = ParagraphStyle(
        'TableBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=11,
        textColor=colors.HexColor("#1E293B")
    )
    
    table_body_bold = ParagraphStyle(
        'TableBodyBold',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9,
        leading=11,
        textColor=colors.HexColor("#1E293B")
    )

    story = []
    
    # --- COVER PAGE ---
    story.append(Spacer(1, 120))
    story.append(Paragraph("SWARNA VAIDEHI MAKHANA HUB", title_style))
    story.append(Paragraph("5-Year Scaling Roadmap, Business Model & Bank Loan DPR", subtitle_style))
    story.append(Spacer(1, 100))
    
    meta_text = """
    <b>Prepared Under the Strategic Framework of:</b><br/>
    ICAR-National Research Centre for Makhana, Darbhanga<br/>
    GI-Tagged Agricultural Crop Expansion Initiative (Bihar, India)<br/><br/>
    <b>Initial Operational Phase:</b> Outsourced/Decentralized Manual Popping Clusters<br/>
    <b>Secondary Operational Phase:</b> 5-Ton Daily Capacity Fully Mechanized Processing Line<br/><br/>
    <b>Date of Issue:</b> May 2026<br/>
    <b>Classification:</b> strictly confidential / bank loan proposal evaluation
    """
    story.append(Paragraph(meta_text, meta_style))
    story.append(PageBreak())
    
    # --- SECTION 1 ---
    story.append(Paragraph("1. Scientific Cultivation Protocol & Value-Add Guide", h1_style))
    story.append(Paragraph("The elite variety <b>'Swarna Vaidehi'</b> (first released on 15.11.2013 and notified on 30.07.2022 by CVRC) developed by the ICAR-NRC for Makhana, Darbhanga represents a massive leap forward. Optimized for field-based agricultural ecosystems (as opposed to natural wetland ponds), this variety ensures high yield stability, standard sizing, and safe manual farming operations.", body_style))
    
    story.append(Paragraph("Step-by-Step Cultivation Protocol:", h2_style))
    story.append(Paragraph("• <b>Nursery Bed Sowing (Nov-March):</b> Raw seeds broadcasted at a low seed rate of <b>20 kg/ha</b> (ponds require 90 kg/ha) in clayey mud nursery beds with an area of 500 m²/ha. Seedlings are ready at 75 days old in February.", bullet_style))
    story.append(Paragraph("• <b>Field Transplanting (Feb-March):</b> Healthy seedlings are transplanted into pre-bunded ploughed fields at a standard row spacing of <b>1.20 m</b>.", bullet_style))
    story.append(Paragraph("• <b>Water Depth Control (March-August):</b> Water is maintained at a shallow depth of strictly <b>30 - 45 cm</b>. This saves 80% of water compared to natural deep ponds and guarantees workers can stand during harvest, bypassing the dangerous deep diving associated with pond monoculture.", bullet_style))
    story.append(Paragraph("• <b>Precision Nutrient Regimen:</b> Soil nutrition requires NPK applied at <b>100:60:40 kg/ha</b> in split doses along with 10 t/ha Farm Yard Manure (FYM). Core secondary micro-nutrients: Zinc (5 kg/ha), Magnesium (5 kg/ha), and Boron (5 kg/ha) combined with foliar spray of 0.1% CuSO4 solution ensure heavy seed shell development.", bullet_style))
    story.append(Paragraph("• <b>Integrated Pest Management (IPM):</b> Sucking insects like Makhana Aphids (*Rhopalosiphum nymphaeae* L.) and Singhara Beetle leaf-perforators are controlled organically using Coccinellid beetles, Metarhizium anisopliae or Beauveria bassiana mycological agents (@ 1x10^8 cfu/g), or 1% Neem oil. Target chemistry includes Flonicamid 50% WG @ 0.4 g/L or Chlorantraniliprole 18.5 SC @ 0.3 mL/L.", bullet_style))
    story.append(Paragraph("• <b>Root Borer Grub Control:</b> Managed via granular Chlorantraniliprole 0.4% GR @ 6 kg/acre applied to standing water. Leaf blight is controlled with Dithane M-45 (0.3%) sprays.", bullet_style))
    story.append(Paragraph("• <b>Harvesting & Sweeping (July-August):</b> Mature seeds (berries grow 6-8 cm, plants bear 8-15 fruits containing 20-200 seeds of 0.8-1.5 cm size) collected standing up using Gaanj nets. Yield ranges from <b>3.0 - 3.5 t/ha</b>.", bullet_style))
    
    story.append(Spacer(1, 10))
    story.append(Paragraph("Premium Value-Added Food Product Chemistry:", h2_style))
    story.append(Paragraph("• <b>Makhana Ragi Cookies (Biscuits):</b> Micronized Euryale ferox seed perisperm flour and finger millet (ragi), rich in essential amino acids (Glutamine, Arginine, Methionine) and calcium (160mg/100g) with stable low fats (0.1%). Retains <b>58% Gross Margin</b>.", bullet_style))
    story.append(Paragraph("• <b>Hypoallergenic Baby Food (Shishu Aahar):</b> Hypoallergenic weaning flour blending roasted Makhana perisperm (45%), decorticated green gram (30%), and ragi (25%). High essential amino acids index (88.5-91.6%) and optimized 1.2:1 Ca:P ratio. Retains <b>62% Gross Margin</b>.", bullet_style))
    story.append(Paragraph("• <b>Low-GI Functional Pasta:</b> Extruded pasta made from 50% cold-gelatinized Euryale starch (high resistant amylose starch) and durum wheat semolina. Highly digestible perisperm starchy core keeping GI below 50. Retains <b>52% Gross Margin</b>.", bullet_style))
    
    story.append(PageBreak())
    
    # --- SECTION 2 ---
    story.append(Paragraph("2. Distributed Sourcing & Logistics Pipeline", h1_style))
    story.append(Paragraph("To capture maximum value-addition and avoid localized middleman manipulation, the Hub structures direct farm-gate buybacks across North Bihar districts (Darbhanga, Madhubani, Purnia):", body_style))
    story.append(Paragraph("1. <b>Farm-Gate Direct Buyback:</b> Raw seeds are purchased directly from bunded fields using digital moisture testers. Prices are locked based on seed size grades: ₹150/kg for Grade-A (>1.2cm), ₹130/kg for Grade-B (0.9-1.2cm), and ₹100/kg for Grade-C (<0.9cm).", bullet_style))
    story.append(Paragraph("2. <b>Primary Sun-Drying:</b> Seeds are spread over HDPE sheets in dedicated sun yards to dry, reducing moisture from 28% down to 12-14%. This prevents aflatoxins and boosts popping ratios.", bullet_style))
    story.append(Paragraph("3. <b>Rotary Nested Sizing:</b> Sorting into 5 distinct sizing drums. Uniform grades are critical to guarantee even heat distribution during roasting.", bullet_style))
    story.append(Paragraph("4. <b>Decentralized Popping Networks:</b> Graded seeds are distributed to skilled traditional Mallah families for popping. Families use the 'double-roast and mallet' method, receiving ₹35/kg job-work fees. This keeps initial CAPEX at zero while creating rural livelihood opportunities.", bullet_style))
    
    # --- SECTION 3 ---
    story.append(Spacer(1, 10))
    story.append(Paragraph("3. Financial Detailed Project Report (DPR)", h1_style))
    story.append(Paragraph("<b>Table 3.1: 1-Hectare Field Cultivation Sheet (Swarna Vaidehi)</b>", h2_style))
    
    t1_data = [
        [Paragraph("<b>Item Description / Input Spec</b>", table_header_style), Paragraph("<b>Cost Category</b>", table_header_style), Paragraph("<b>Cost (INR)</b>", table_header_style)],
        [Paragraph("Land Lease of bunded field (1 Year)", table_body_style), Paragraph("Fixed Lease Cost", table_body_style), Paragraph("Rs. 25,000", table_body_bold)],
        [Paragraph("Field preparation & perimeter leveling", table_body_style), Paragraph("Tractor & Earth Moving", table_body_style), Paragraph("Rs. 15,000", table_body_bold)],
        [Paragraph("Nursery seedlings & certified seed stock", table_body_style), Paragraph("Seeds (20 kg/ha)", table_body_style), Paragraph("Rs. 8,000", table_body_bold)],
        [Paragraph("NPK Fertilizer & micro-nutrient nutrition", table_body_style), Paragraph("Soil Inputs (Zn, Mg, B)", table_body_style), Paragraph("Rs. 12,000", table_body_bold)],
        [Paragraph("Integrated Pest Management bio-agents", table_body_style), Paragraph("Crop Protection", table_body_style), Paragraph("Rs. 5,000", table_body_bold)],
        [Paragraph("Transplanting, weeding & sweeping labor", table_body_style), Paragraph("75 Mandays @ Rs. 600", table_body_style), Paragraph("Rs. 45,000", table_body_bold)],
        [Paragraph("Irrigation pump diesel fuel", table_body_style), Paragraph("Water management", table_body_style), Paragraph("Rs. 10,000", table_body_bold)],
        [Paragraph("<b>TOTAL CULTIVATION COST (OPEX)</b>", table_body_bold), Paragraph("<b>Total Operational Expense</b>", table_body_bold), Paragraph("<b>Rs. 1,20,000</b>", table_body_bold)],
        [Paragraph("Raw Seed Revenue (2800 kg @ Rs. 150/kg)", table_body_style), Paragraph("Market Seed Sales", table_body_style), Paragraph("Rs. 4,20,000", table_body_bold)],
        [Paragraph("<b>NET PROFIT PER HECTARE</b>", table_body_bold), Paragraph("<b>Net Cultivator Return</b>", table_body_bold), Paragraph("<b>Rs. 3,00,000</b>", table_body_bold)],
    ]
    
    t1 = Table(t1_data, colWidths=[200, 160, 140])
    t1.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1B4332")),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#CBD5E1")),
        ('ROWBACKGROUNDS', (0,1), (-1,-3), [colors.white, colors.HexColor("#F8F9FA")]),
        ('BACKGROUND', (0,8), (-1,8), colors.HexColor("#D8F3DC")),
        ('BACKGROUND', (0,10), (-1,10), colors.HexColor("#FEFAE0")),
        ('PADDING', (0,0), (-1,-1), 5),
    ]))
    story.append(t1)
    
    story.append(PageBreak())
    
    # --- PROCESSING LEDGER ---
    story.append(Paragraph("<b>Table 3.2: Year-1 Bootstrapped Processing Ledger (30 Tons seeds)</b>", h2_style))
    story.append(Paragraph("Operational setup with <i>zero machinery purchase</i>, utilizing outsourced traditional Mallah popping clusters for double roasting and mallet popping.", body_style))
    
    t2_data = [
        [Paragraph("<b>Operational Component</b>", table_header_style), Paragraph("<b>Specification & Scale</b>", table_header_style), Paragraph("<b>Cost / Value (INR)</b>", table_header_style)],
        [Paragraph("Raw Seed Procurement", table_body_style), Paragraph("30,000 kg seeds @ average Rs. 130/kg", table_body_style), Paragraph("Rs. 39,00,000", table_body_bold)],
        [Paragraph("Drying & Sorting Labor", table_body_style), Paragraph("150 Mandays @ Rs. 800 at hubs", table_body_style), Paragraph("Rs. 1,20,000", table_body_bold)],
        [Paragraph("Logistics & Transport", table_body_style), Paragraph("Farm-gate to processing clusters", table_body_style), Paragraph("Rs. 1,50,000", table_body_bold)],
        [Paragraph("Manual Popping Job-work", table_body_style), Paragraph("11,700 kg popped makhana @ Rs. 35/kg", table_body_style), Paragraph("Rs. 4,09,500", table_body_bold)],
        [Paragraph("Branding & Pouch Packaging", table_body_style), Paragraph("46,800 custom polymer bags", table_body_style), Paragraph("Rs. 1,20,000", table_body_bold)],
        [Paragraph("Warehouse Rent & Overheads", table_body_style), Paragraph("12 months general hub costs", table_body_style), Paragraph("Rs. 1,50,000", table_body_bold)],
        [Paragraph("<b>TOTAL SOURCING & PROCESSING OPEX</b>", table_body_bold), Paragraph("<b>Total Sourced Cost</b>", table_body_bold), Paragraph("<b>Rs. 48,49,500</b>", table_body_bold)],
        [Paragraph("Revenue: Premium Popped Grade-A", table_body_style), Paragraph("8,000 kg sold wholesale @ Rs. 550/kg", table_body_style), Paragraph("Rs. 44,00,000", table_body_bold)],
        [Paragraph("Revenue: Standard Popped Grade-B", table_body_style), Paragraph("3,000 kg sold wholesale @ Rs. 400/kg", table_body_style), Paragraph("Rs. 12,00,000", table_body_bold)],
        [Paragraph("Revenue: Broken / Phulla Grade-C", table_body_style), Paragraph("700 kg sold for weaning flour @ Rs. 200/kg", table_body_style), Paragraph("Rs. 1,40,000", table_body_bold)],
        [Paragraph("Revenue: Makhana Shell Waste", table_body_style), Paragraph("18,000 kg sold as mulch @ Rs. 3/kg", table_body_style), Paragraph("Rs. 54,000", table_body_bold)],
        [Paragraph("<b>TOTAL REVENUE FROM POPPING</b>", table_body_bold), Paragraph("<b>Total Gross Returns</b>", table_body_bold), Paragraph("<b>Rs. 57,94,000</b>", table_body_bold)],
        [Paragraph("<b>NET HUB PROCESSING PROFIT (Y1)</b>", table_body_bold), Paragraph("<b>Net Processing Hub Profit</b>", table_body_bold), Paragraph("<b>Rs. 9,44,500</b>", table_body_bold)],
    ]
    
    t2 = Table(t2_data, colWidths=[180, 180, 140])
    t2.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1B4332")),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#CBD5E1")),
        ('ROWBACKGROUNDS', (0,1), (-1,-8), [colors.white, colors.HexColor("#F8F9FA")]),
        ('BACKGROUND', (0,7), (-1,7), colors.HexColor("#D8F3DC")),
        ('BACKGROUND', (0,12), (-1,12), colors.HexColor("#D8F3DC")),
        ('BACKGROUND', (0,13), (-1,13), colors.HexColor("#FEFAE0")),
        ('PADDING', (0,0), (-1,-1), 4),
    ]))
    story.append(t2)
    
    story.append(Spacer(1, 10))
    story.append(Paragraph("<b>Consolidated Profit Synthesis:</b> By integrating 1-Hectare self-cultivation (₹3,00,000) and the 30-ton hub ledger (₹9,44,500), the business achieves an aggregate net profit of <b>₹12,44,500 in Year 1</b> with zero machinery CAPEX. This fulfills the strategic objective.", body_style))
    
    # --- FUTURE 5-TON SCALING PHASE ---
    story.append(Paragraph("<b>Table 3.3: Future 5-Ton Mechanized Scaling Phase (CAPEX Breakdown)</b>", h2_style))
    story.append(Paragraph("CAPEX values are strictly derived from the **ICAR-NRC Makhana core food machinery price records** (File 1):", body_style))
    
    t3_data = [
        [Paragraph("<b>Machinery CAPEX Items</b>", table_header_style), Paragraph("<b>Technical Details / Hourly Output</b>", table_header_style), Paragraph("<b>Cost (INR)</b>", table_header_style)],
        [Paragraph("Secondary Popping Machine", table_body_style), Paragraph("Impact-based popping at 270 °C (25-30 kg/h)", table_body_style), Paragraph("Rs. 8,00,000", table_body_bold)],
        [Paragraph("Primary Seed Roasting Machine", table_body_style), Paragraph("Roasts raw seeds at 270 °C (100 kg/h)", table_body_style), Paragraph("Rs. 1,40,000", table_body_bold)],
        [Paragraph("Rotary Nested Seed Grader", table_body_style), Paragraph("Separates raw seeds into 7 grades (180 kg/h)", table_body_style), Paragraph("Rs. 2,50,000", table_body_bold)],
        [Paragraph("Makhana Seed Washer", table_body_style), Paragraph("Cleans raw seeds via friction (120-180 kg/h)", table_body_style), Paragraph("Rs. 1,60,000", table_body_bold)],
        [Paragraph("Makhana Seed Grader", table_body_style), Paragraph("Sand-suction slurry sorting (0.5 t/h)", table_body_style), Paragraph("Rs. 1,30,000", table_body_bold)],
        [Paragraph("Cabinet Seed Dryer", table_body_style), Paragraph("Temp range 30-220 °C (20-30 kg/h)", table_body_style), Paragraph("Rs. 2,00,000", table_body_bold)],
        [Paragraph("Popped Makhana Grader", table_body_style), Paragraph("Sorts puffs into 7 size classes (100 kg/h)", table_body_style), Paragraph("Rs. 2,00,000", table_body_bold)],
        [Paragraph("Automatic FFS Packing Machine", table_body_style), Paragraph("Multi-head weigher packing (6-8 t/day)", table_body_style), Paragraph("Rs. 8,00,000", table_body_bold)],
        [Paragraph("Makhana Polishing Unit", table_body_style), Paragraph("Polishes finished pops to standard textures", table_body_style), Paragraph("Rs. 1,00,000", table_body_bold)],
        [Paragraph("Warehouse Civil Works & Drying Slab", table_body_style), Paragraph("Concrete slab and dust-free sorting floor", table_body_style), Paragraph("Rs. 15,00,000", table_body_bold)],
        [Paragraph("Solar Rooftop Micro-Grid", table_body_style), Paragraph("15kW solar array with 10kW battery backup panels", table_body_style), Paragraph("Rs. 5,00,000", table_body_bold)],
        [Paragraph("<b>TOTAL CAPEX REQUIREMENT</b>", table_body_bold), Paragraph("<b>Automated Infrastructure CAPEX</b>", table_body_bold), Paragraph("<b>Rs. 48,50,000</b>", table_body_bold)],
    ]
    
    t3 = Table(t3_data, colWidths=[160, 200, 140])
    t3.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1B4332")),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#CBD5E1")),
        ('ROWBACKGROUNDS', (0,1), (-1,-2), [colors.white, colors.HexColor("#F8F9FA")]),
        ('BACKGROUND', (0,12), (-1,12), colors.HexColor("#D8F3DC")),
        ('PADDING', (0,0), (-1,-1), 4),
    ]))
    story.append(t3)
    
    story.append(PageBreak())
    
    # --- SECTION 4 ---
    story.append(Paragraph("4. Integrated Farming System (IFS) & Excellence Analysis", h1_style))
    story.append(Paragraph("To maximize field efficiency, the traditional makhana monoculture is compared against an **Integrated Farming System (IFS)**. Stocking Indian Major Carps (Catla, Rohu, Mrigal) and planting a winter crop of Water Chestnut (Singhara) in the same field body increases resource efficiency.", body_style))
    
    t4_data = [
        [Paragraph("<b>Operational Metric</b>", table_header_style), Paragraph("<b>Traditional Field Monoculture</b>", table_header_style), Paragraph("<b>Makhana + Carps + Water Chestnut IFS</b>", table_header_style)],
        [Paragraph("Water Depth Maintenance", table_body_style), Paragraph("30 - 45 cm (Shallow)", table_body_style), Paragraph("50 - 60 cm (For carp movement)", table_body_bold)],
        [Paragraph("Primary Makhana Yield", table_body_style), Paragraph("2,800 kg/ha", table_body_style), Paragraph("2,600 kg/ha (Due to fish channels)", table_body_bold)],
        [Paragraph("Carps Fish Yield", table_body_style), Paragraph("0 kg (None)", table_body_style), Paragraph("1,200 kg/ha (Synergistic carp growth)", table_body_bold)],
        [Paragraph("Water Chestnut Yield", table_body_style), Paragraph("0 kg (None)", table_body_style), Paragraph("3,000 kg/ha (Secondary winter crop)", table_body_bold)],
        [Paragraph("Total Operational Cost", table_body_style), Paragraph("Rs. 1,20,000", table_body_style), Paragraph("Rs. 1,80,000 (Includes fingerlings/seeds)", table_body_bold)],
        [Paragraph("Total Gross Revenue", table_body_style), Paragraph("Rs. 4,20,000", table_body_style), Paragraph("Rs. 7,30,000 (Makhana + Fish + Chestnut)", table_body_bold)],
        [Paragraph("<b>NET FARM PROFIT PER HA</b>", table_body_bold), Paragraph("<b>Rs. 3,00,000</b>", table_body_bold), Paragraph("<b>Rs. 5,50,000 (An 83% Profit Increase)</b>", table_body_bold)],
    ]
    
    t4 = Table(t4_data, colWidths=[180, 160, 160])
    t4.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1B4332")),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#CBD5E1")),
        ('ROWBACKGROUNDS', (0,1), (-1,-2), [colors.white, colors.HexColor("#F8F9FA")]),
        ('BACKGROUND', (0,7), (-1,7), colors.HexColor("#FEFAE0")),
        ('PADDING', (0,0), (-1,-1), 5),
    ]))
    story.append(t4)
    
    story.append(Spacer(1, 15))
    story.append(Paragraph("Organic Nutrient & Pest Protection Sourcing:", h2_style))
    story.append(Paragraph("• <b>Azolla Pinata:</b> Grown floating in early makhana vegetative phases, this bio-fertilizer fixes nitrogen, reducing urea dependency by <b>40%</b>.", bullet_style))
    story.append(Paragraph("• <b>Neem Cake (250 kg/ha):</b> Basal dressing during soil prep, naturally suppressing root-knot nematodes and eliminating synthetic chemistry.", bullet_style))
    story.append(Paragraph("• <b>Garlic-Chilli botanical sprays:</b> Made locally, protecting the early foliage against case-worms (*Nymphula crisonalis*) and aphids without chemical residues.", bullet_style))
    story.append(Paragraph("• <b>Trichoderma Viride:</b> Fungal treatment of nursery seed stock to prevent damping-off.", bullet_style))
    
    story.append(Spacer(1, 10))
    story.append(Paragraph("Global Export Positioning:", h2_style))
    story.append(Paragraph("• <b>United States (US):</b> Organic puffed snacking segment, fetching an **85% price premium**.", bullet_style))
    story.append(Paragraph("• <b>European Union (EU):</b> Rigid multi-residue pesticide limits and high aflatoxin compliance (< 2 ppb), yielding a **110% price premium**.", bullet_style))
    story.append(Paragraph("• <b>Bihar Bazar Integration:</b> Live integration with `bihar-bazaar.vercel.app` enables voice/SMS price discovery, bypassing local middlemen and providing transparent batch tracking for international buyers.", bullet_style))
    
    doc.build(story, canvasmaker=NumberedCanvas)
    print("PDF Report saved successfully!")

# -------------------------------------------------------------
# MAIN RUNNER
# -------------------------------------------------------------
if __name__ == "__main__":
    print("Starting generation of all document formats...")
    generate_excel()
    generate_ppt()
    generate_pdf()
    print("All file formats generated successfully!")
